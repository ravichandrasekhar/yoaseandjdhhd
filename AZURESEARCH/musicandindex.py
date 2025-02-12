import requests
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    HnswAlgorithmConfiguration,
    VectorSearch,
    VectorSearchAlgorithmConfiguration,
    VectorSearchProfile,
)
import fitz
from azure.search.documents.indexes.models import (
    SearchIndex,
    SimpleField,
    SearchableField,
    SearchField
)
from pptx import Presentation
import pdfplumber
from docx import Document
import io
import os
import csv
import openpyxl
import numpy as np
from azure.ai.textanalytics import TextAnalyticsClient
from azure.core.credentials import AzureKeyCredential
from openai import AzureOpenAI
from azure.identity import DefaultAzureCredential, get_bearer_token_provider
import tempfile
from langchain_community.document_loaders import PyPDFLoader
import re
import PyPDF2
import logging
from bs4 import BeautifulSoup
folder_id = ""
SHAREPOINT_DRIVE_ID = ""
ACCESS_TOKEN = ""
SEARCH_ENDPOINT = ""
SEARCH_ADMIN_KEY = ""
SEARCH_INDEX_NAME = ""
group_id = ""
azure_openai_endpoint = ""
azure_openai_key = ""
azure_openai_embedding_deployment = "text-embedding-ada-002"
embedding_model_name = "text-embedding-ada-002"
azure_openai_version="2023-05-15"
def create_index():
    index_fields = [
        SimpleField(name="id", type="Edm.String", key=True,filterable=True, sortable=True),
        SearchableField(name="file_name", type="Edm.String",filterable=True),
        SearchableField(name="title", type="Edm.String"),
        SearchableField(name="content", type="Edm.String"),
        SimpleField(name="folderName", type="Edm.String"),
        SimpleField(name="webUrl", type="Edm.String"),
        SearchField(name="contentVector", type="Collection(Edm.Single)",vector_search_dimensions=1536, vector_search_profile_name="myHnswProfile")
        # Add more fields as needed
    ]
    vector_search = VectorSearch(
        algorithms=[
            HnswAlgorithmConfiguration(
                name="myHnsw"
            )
        ],
        profiles=[
            VectorSearchProfile(
                name="myHnswProfile",
                algorithm_configuration_name="myHnsw",
            )
        ]
    )
    index = SearchIndex(name=SEARCH_INDEX_NAME, fields=index_fields,vector_search=vector_search)

    # Create index client
    index_client = SearchIndexClient(SEARCH_ENDPOINT, AzureKeyCredential(SEARCH_ADMIN_KEY))

    try:
        # Check if the index exists
        existing_index = index_client.get_index(SEARCH_INDEX_NAME)
        print(f"Index '{SEARCH_INDEX_NAME}' already exists. Deleting the existing index...")
        index_client.delete_index(SEARCH_INDEX_NAME)
        print(f"Index '{SEARCH_INDEX_NAME}' has been deleted.")
    except Exception as e:
        print(f"Index '{SEARCH_INDEX_NAME}' not found.")

    try:
        print(f"Creating index '{SEARCH_INDEX_NAME}'...")
        index_client.create_index(index)
        print(f"Index '{SEARCH_INDEX_NAME}' has been created successfully.")
    except Exception as e:
        print(f"Error creating index '{SEARCH_INDEX_NAME}': {e}")

def get_sharepoint_items(access_token, folder_id=None):
    
    if folder_id:
        url = f"https://graph.microsoft.com/v1.0/drives/{SHAREPOINT_DRIVE_ID}/items/{folder_id}/children"
    else:
        url = f"https://graph.microsoft.com/v1.0/drives/{SHAREPOINT_DRIVE_ID}/root/children"

    headers = {
        "Authorization": "Bearer " + access_token,
        "Accept": "application/json"
    }
    response = requests.get(url, headers=headers)
    data = response.json()
    return data['value']


def download_sharepoint_file(file_id, access_token):
    content_url = f"https://graph.microsoft.com/v1.0/drives/{SHAREPOINT_DRIVE_ID}/items/{file_id}/content"
    headers = {
        "Authorization": "Bearer " + access_token,
    }
    response = requests.get(content_url, headers=headers)
    if response.status_code == 200:
        return response.content
    else:
        print(f"Failed to download file with ID '{file_id}'")
        return None

def convert_to_html(content_str):
    """
    Converts plain text content into HTML using BeautifulSoup.
    """

    # Initialize BeautifulSoup object
    soup = BeautifulSoup("<html><body></body></html>", "html.parser")

    # Split the content into paragraphs, assuming '\n' is a paragraph separator
    paragraphs = content_str.split("\n")
    
    for paragraph in paragraphs:
        if paragraph.strip():  # Ignore empty paragraphs
            # Create a <p> tag for each paragraph and append to body
            p_tag = soup.new_tag("p")
            p_tag.string = paragraph.strip()  # Clean up leading/trailing spaces
            soup.body.append(p_tag)

    # Return the prettified HTML content
    return soup.prettify()
def get_permissions(file_id, access_token):
    permissions_url = f"https://graph.microsoft.com/v1.0/drives/{SHAREPOINT_DRIVE_ID}/items/{file_id}/permissions"
    headers = {
        "Authorization": "Bearer " + access_token,
        "Accept": "application/json"
    }
    response = requests.get(permissions_url, headers=headers)
    permissions = []
    if response.status_code == 200:
        permissions_data = response.json().get("value", [])
        for permission_data in permissions_data:
            granted_to = permission_data.get("grantedTo", None)
            if granted_to:
                if "user" in granted_to:
                    user_or_group = granted_to["user"].get("displayName", "N/A")
                elif "group" in granted_to:
                    user_or_group = granted_to["group"].get("displayName", "N/A")
                if user_or_group != "N/A":  # Filter out "N/A" values
                    permissions.append(f" {user_or_group}")

    return permissions
def process_pdf(pdf_content):
    doc = fitz.open(stream=io.BytesIO(pdf_content)) 
    full_text = [] 
    for page in doc: 
        page_content = page.get_text() 
        full_text.append(page_content) 
    return full_text
 


def get_group_members(group_id, access_token):
    url = f"https://graph.microsoft.com/v1.0/groups/{group_id}/transitiveMembers/"
    group_members = []
    headers = {
        "Authorization": "Bearer " + access_token,
        "Accept": "application/json"
    }

    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        members = response.json()["value"]

        for member in members:
            principal_name = member["userPrincipalName"]
            group_members.append(principal_name)

    else:
        print(f"Failed to retrieve the group members: {response.text}")

    return group_members

def process_pdf(pdf_content):
        """Extract text from a PDF file and remove page numbers."""
        try:
            doc = fitz.open(stream=io.BytesIO(pdf_content))
            page_text = {}
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                clean_text = [line for line in text.splitlines() if not re.match(r"(Page \d+|Page \d+ of \d+)", line.strip())]
                page_text[f"Page {page_num}"] = "\n".join(clean_text)
        except Exception as e:
            logging.error(f"Error processing PDF file: {str(e)}")
            return {}
        return page_text
def process_pdf(pdf_content):
    """Extracts text from PDF and returns it as a list of lines."""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_content))
    chunks = []
    
    # Iterate through each page in the PDF
    for page in pdf_reader.pages:
        page_content = page.extract_text()
        
        if page_content:  # Check if there is content on the page
            # Split the page content into lines based on single newlines
            lines = page_content.split('\n')
            
            # Add lines to chunks, avoiding empty lines
            for line in lines:
                if line.strip():
                    chunks.append(line.strip())
    
    return chunks

def extract_text_from_pdf(file_content):
    """Extracts text from PDF and returns it as a list of lines."""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
    chunks = []
        
        # Iterate through each page in the PDF
    for page in pdf_reader.pages:
            page_content = page.extract_text()
            
            if page_content:  # Check if there is content on the page
                # Split the page content into paragraphs based on double newlines
                paragraphs = page_content.split('\n\n')
                
                # Add paragraphs to chunks, avoiding empty paragraphs
                for paragraph in paragraphs:
                    if paragraph.strip():
                        chunks.append(paragraph.strip())
        
    return " ".join(chunks)

def extract_text_from_word(file_content):
    doc = Document(io.BytesIO(file_content))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text
def extract_text_from_pptx(file_content):
    """
    Extracts text from a PowerPoint (.pptx) file.
    """
    try:
        presentation = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
            text += "\n"  # Separate slides with a newline
        return text.strip()
    except Exception as e:
        return f"Error extracting text from PowerPoint file: {str(e)}"

def extract_text_from_excel(file_content):
    file_str = ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_content))
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        file_str += str(cell.value) + " "
    except Exception as e:
        file_str = f"Error extracting text from Excel file: {str(e)}"
    return file_str

def extract_text_from_text(file_content):
    try:
        return file_content.decode('utf-8')
    except UnicodeDecodeError:
        return str(file_content, errors='ignore')


def extract_text_from_other(file_content):
    # Placeholder for handling other file types
    return "Unsupported file type"


def extract_text_from_csv(file_content):
    file_str = ""
    try:
        # Decode the bytes content assuming utf-8 encoding
        content_str = file_content.decode('utf-8')
        # Parse CSV content
        reader = csv.reader(content_str.splitlines())
        for row in reader:
            for cell in row:
                file_str += cell + " "
    except Exception as e:
        file_str = f"Error extracting text from CSV file: {str(e)}"
    return file_str

def index_folder_to_azure_search(folder_name, folder_id):
    client = SearchClient(endpoint=SEARCH_ENDPOINT,
                          index_name=SEARCH_INDEX_NAME,
                          credential=AzureKeyCredential(SEARCH_ADMIN_KEY))

    document = {
        "id": str(folder_id),  # Unique identifier for the folder
        "folderName": folder_name
        # Add more fields as needed
    }

    try:
        result = client.upload_documents(documents=[document])
        print(f"Folder '{folder_name}' indexed successfully.")
    except Exception as e:
        print(
            f"Failed to index folder '{folder_name}' to Azure Search. Error: {str(e)}")
def index_document_to_azure_search(file_name, file_content, file_id, web_url):
    file_extension = os.path.splitext(file_name)[1].lower()
    search_client = SearchClient(endpoint=SEARCH_ENDPOINT,
                              index_name=SEARCH_INDEX_NAME,
                              credential=AzureKeyCredential(SEARCH_ADMIN_KEY))
    content_str = ""

    try:
        if file_extension == ".docx" or file_extension == ".doc":
            content_str = extract_text_from_word(file_content)
        elif file_extension == ".pdf":
            content_str = extract_text_from_pdf(file_content)
            print("content from pdf",content_str)
        elif file_extension == ".txt":
            content_str = extract_text_from_text(file_content)
        elif file_extension == ".xlsx" or file_extension == ".xls":
            content_str = extract_text_from_excel(file_content)
        elif file_extension == ".csv":
            content_str = extract_text_from_csv(file_content)
        elif file_extension == ".pptx":
            content_str = extract_text_from_pptx(file_content)
            print("content from pptx",content_str)
        else:
            content_str = extract_text_from_other(file_content)

    except Exception as e:
        print(f"Error extracting text from '{file_name}': {str(e)}")
        return
    content_str = str(content_str).strip()
    # html_content = convert_to_html(content_str)
    # print("html content",html_content)
    chunk_size = 8000  # Define chunk size as needed
    # content_chunks = [content_str[i:i + chunk_size] for i in range(0, len(content_str), chunk_size)]
    content_chunks = [content_str[i:i + chunk_size] for i in range(0, len(content_str), chunk_size)]
    try:
        # Azure OpenAI client initialization and content embedding generation
        openai_credential = DefaultAzureCredential()
        token_provider = get_bearer_token_provider(openai_credential, "https://cognitiveservices.azure.com/.default")
        client = AzureOpenAI(
            azure_deployment=azure_openai_embedding_deployment,
            azure_endpoint=azure_openai_endpoint,
            api_key=azure_openai_key,
            azure_ad_token_provider=token_provider if not azure_openai_key else None,
            api_version=azure_openai_version
        )
        embeddings=[]
        for content in content_chunks:
        
                    try:
                        # Generate embeddings for content chunk
                        response = client.embeddings.create(input=content, model=embedding_model_name)
                        if response.data and response.data[0].embedding:  # Check if response and embedding exist
                            # Extract embeddings and ensure the correct dimensionality
                            embeddings_chunk = response.data[0].embedding[:1536]  # Limit to 1536 dimensions
                            embeddings.append(embeddings_chunk)
                        else:
                            print("Error: Embedding data not found in response.")
                    except Exception as e:
                        print(f"Error generating embeddings for content chunk: {e}")

                # Convert embeddings to numpy array
        embeddings_array = np.array(embeddings)

                # Convert embeddings array to list for JSON serialization
        embeddings_list = embeddings_array.tolist()


        document = {
            "id": str(file_id),  # Unique identifier for the document
            "file_name": file_name,
            "title": file_name,
            "content": content_str,
            "webUrl": web_url,
            "contentVector": embeddings_list[0] if embeddings_list else []
        }
        print
        result = search_client.upload_documents(documents=[document])
        print(f"Document '{file_name}' indexed successfully.")

    
    except Exception as e:
        print(f"Error indexing document '{file_name}' to Azure Search: {e}")



def index_sharepoint_items(access_token, parent_folder_id=None):
    items = get_sharepoint_items(access_token, parent_folder_id)
    for item in items:
        item_name = item['name']
        item_id = item['id']
        web_url = item['webUrl']
        if 'folder' in item:
            index_folder_to_azure_search(item_name, item_id)
            index_sharepoint_items(access_token, item_id)
        else:
            file_content = download_sharepoint_file(item_id, access_token)
            
            if file_content:
                # Extract text based on file type
                if item_name.endswith('.pdf'):
                    content_str = extract_text_from_pdf(file_content)
                elif item_name.endswith('.docx'):
                    content_str = extract_text_from_word(file_content)
                elif item_name.endswith('.xlsx') or item_name.endswith('.xls'):
                    content_str = extract_text_from_excel(file_content)
                else:
                    content_str = extract_text_from_other(file_content)
                
                # Generate content embeddings using Azure OpenAI
                index_document_to_azure_search(item_name, file_content, item_id, web_url)


def main():
    global client
    
    openai_credential = DefaultAzureCredential()
    token_provider = get_bearer_token_provider(
        openai_credential, "https://cognitiveservices.azure.com/.default")
    
    client = AzureOpenAI(
        azure_deployment=azure_openai_embedding_deployment,
        azure_endpoint=azure_openai_endpoint,
        api_key=azure_openai_key,
        azure_ad_token_provider=token_provider if not azure_openai_key else None,
        api_version=azure_openai_version
    )
    
    create_index()
    index_sharepoint_items(ACCESS_TOKEN, parent_folder_id=folder_id)


if __name__ == "__main__":
    main()