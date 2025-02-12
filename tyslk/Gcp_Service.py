import os
import pandas as pd
from google.cloud import documentai_v1 as documentai
from docx import Document as DocxDocument
from docx2txt import process as extract_images_from_docx
from pptx import Presentation
from dotenv import load_dotenv
from io import BytesIO
from typing import Dict, Any
from Services.Text_Extraction.iText_Extraction import Text_Service

class GcpService(Text_Service):
    def __init__(self):
        load_dotenv()

        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = os.path.join(os.getcwd(), "credentials.json")

        self.project_id = os.getenv("GCP_PROJECT_ID")
        self.location = os.getenv("GCP_LOCATION")
        self.processor_id = os.getenv("GCP_PROCESSOR_ID")

        # Initialize GCP clients
        self.documentai_client = documentai.DocumentProcessorServiceClient()

    def validate_config(self) -> Dict[str, Any]:
        required_fields = ["GCP_PROJECT_ID", "GCP_LOCATION", "GCP_PROCESSOR_ID"]
        missing_fields = [fields for fields in required_fields if not os.getenv(fields)]

        if missing_fields:
            return {
                "status": "error",
                "message": f"Missing required environment variables: {', '.join(missing_fields)}",
                "error": True
            }

        credentials_path = os.path.join(os.getcwd(), "credentials.json")
        if not os.path.isfile(credentials_path):
            return {
                "status": "error",
                "message": "credentials.json file is missing",
                "error": True
            }

        return {
            "status": "success",
            "message": "Initialization config validated successfully",
            "error": False
        }

    def _extract_text_from_pdf_or_image(self, object_content, mime_type):
        """Extracts text and metadata from PDF or image files using GCP Document AI or Vision API."""
        if mime_type == "application/pdf":
            # Using Document AI for PDFs
            document = documentai.RawDocument(content=object_content, mime_type=mime_type)
            request = documentai.ProcessRequest(
                name=f"projects/{self.project_id}/locations/{self.location}/processors/{self.processor_id}",
                raw_document=document
            )
            response = self.documentai_client.process_document(request=request)
            extracted_text = response.document.text if response.document.text else "No text found."
            metadata = self._extract_metadata_from_document(response.document)
        elif mime_type in ["image/jpeg", "image/png", "image/jpg"]:
            # Using Document AI for images
            document = documentai.RawDocument(content=object_content, mime_type=mime_type)
            request = documentai.ProcessRequest(
                name=f"projects/{self.project_id}/locations/{self.location}/processors/{self.processor_id}",
                raw_document=document
            )
            response = self.documentai_client.process_document(request=request)
            extracted_text = response.document.text if response.document.text else "No text found."
            metadata = self._extract_metadata_from_document(response.document)
        else:
            raise ValueError(f"Unsupported MIME type: {mime_type}")

        return extracted_text, metadata

    def _extract_metadata_from_document(self, document):
        """Extract metadata from the Document AI response."""
        metadata = {
            "pages": len(document.pages),
            "document_type": document.document_type,
            "entities": [],
            "tables": [],
            "forms": []
        }

        # Collect entities (e.g., detected key-value pairs)
        for entity in document.entities:
            metadata["entities"].append({
                "type": entity.type_,
                "mention_text": entity.mention_text,
                "value": entity.value.text if entity.value else None
            })

        # Collect tables and forms (if any)
        for page in document.pages:
            for table in page.tables:
                metadata["tables"].append(table)
            for form in page.forms:
                metadata["forms"].append(form)

        return metadata

    def _extract_text_from_docx_with_images(self, object_content):
        """Extracts text and embedded images from a DOCX file and processes images for text extraction."""
        with open("temp_file.docx", "wb") as f:
            f.write(object_content)
        
        doc = DocxDocument("temp_file.docx")
        text = "\n".join([para.text for para in doc.paragraphs])

        # Extract embedded images from DOCX file
        images_dir = "temp_images"
        os.makedirs(images_dir, exist_ok=True)
        extract_images_from_docx("temp_file.docx", images_dir)

        image_texts = []
        for image_filename in os.listdir(images_dir):
            if image_filename.endswith((".png", ".jpg", ".jpeg")):
                image_path = os.path.join(images_dir, image_filename)
                
                with open(image_path, "rb") as image_file:
                    image_bytes = image_file.read()

                # OCR the image
                image_text = self._extract_text_from_pdf_or_image(image_bytes, mime_type="image/jpeg")
                image_texts.append(image_text)

        # Cleanup temporary files and directory
        os.remove("temp_file.docx")
        for image_filename in os.listdir(images_dir):
            os.remove(os.path.join(images_dir, image_filename))
        os.rmdir(images_dir)

        return text + "\n".join(image_texts), {"image_count": len(image_texts)}

    def _extract_text_from_csv(self, object_content):
        """Extracts text from a CSV file."""
        csv_str = object_content.decode('utf-8')
        df = pd.read_csv(BytesIO(csv_str.encode('utf-8')))
        return df.to_string(), {"row_count": len(df)}

    def _extract_text_from_xlsx(self, object_content):
        """Extracts text from an XLSX file."""
        xlsx_file = BytesIO(object_content)
        df = pd.read_excel(xlsx_file, sheet_name=None)
        sheet_info = [{"sheet_name": sheet, "row_count": len(df[sheet])} for sheet in df]
        return "\n\n".join([f"Sheet: {sheet}\n{df[sheet].to_string()}" for sheet in df]), sheet_info

    def _extract_text_from_txt(self, object_content):
        """Extracts text from a TXT file."""
        return object_content.decode('utf-8'), {}

    def _extract_text_from_pptx(self, object_content):
        """Extracts text from a PPTX file."""
        pptx_file = BytesIO(object_content)
        presentation = Presentation(pptx_file)
        text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return "\n".join(text), {"slide_count": len(presentation.slides)}

    def process(self, object_content, object_name):
        """Extracts text and metadata from a specific object file (byte content)."""
        print(f"Processing object: {object_name}")
        
        # Extract text and metadata from the object based on its type
        if object_name.endswith(".pdf"):
            return self._extract_text_from_pdf_or_image(object_content, mime_type="application/pdf")
        elif object_name.endswith((".png", ".jpg", ".jpeg")):
            return self._extract_text_from_pdf_or_image(object_content, mime_type="image/jpeg")
        elif object_name.endswith(".docx"):
            return self._extract_text_from_docx_with_images(object_content)
        elif object_name.endswith(".csv"):
            return self._extract_text_from_csv(object_content)
        elif object_name.endswith(".xlsx"):
            return self._extract_text_from_xlsx(object_content)
        elif object_name.endswith(".txt"):
            return self._extract_text_from_txt(object_content)
        elif object_name.endswith(".pptx"):
            return self._extract_text_from_pptx(object_content)
        else:
            raise ValueError("Unsupported file format.")

    def process_file(self, object_name, object_content):
        """Main method to process a file and return both extracted text and metadata."""
        extracted_text, metadata = self.process(object_content, object_name)
        
        return {
            "object_name": object_name,
            "extracted_text": extracted_text,
            "metadata": metadata
        }
