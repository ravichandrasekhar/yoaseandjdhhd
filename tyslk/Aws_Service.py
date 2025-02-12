import os
import time
import boto3
import pandas as pd
from io import BytesIO
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv
import re
from docx2txt import process as extract_images_from_docx
from botocore.exceptions import NoCredentialsError
from typing import Dict, Any
from Services.Text_Extraction.iText_Extraction import Text_Service
 
class AwsService(Text_Service):
    def __init__(self):
        load_dotenv()
        self.textract_client = boto3.client('textract', region_name=os.getenv("AWS_REGION"))
       
    def validate_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Validates the AWS configuration for Textract."""
        required_fields = ["AWS_ACCESS_KEY_ID", "AWS_SECRET_ACCESS_KEY", "AWS_REGION"]
        missing_fields = [fields for fields in required_fields if not config['config'].get(fields)]
 
        if missing_fields:
            return {
                "status": "error",
                "message": f"Missing required environment variables for AWS: {', '.join(missing_fields)}",
                "error": True
            }
 
        return {
            "status": "success",
            "message": "AWS Textract configuration validated successfully",
            "error": False
        }
 
    def _process_pdf_or_image_with_textract(self, file_content, metadata):
        """Processes PDF or image files using AWS Textract for text extraction."""
        try:
            response = self.textract_client.analyze_document(
                Document={'Bytes': file_content},
                FeatureTypes=['TABLES', 'FORMS']  # Optional: Include additional features like tables and forms
            )
            extracted_text = ''
            for block in response.get('Blocks', []):
                if block['BlockType'] == 'LINE':
                    extracted_text += block['Text'] + "\n"
            return {"extracted_text": extracted_text, "metadata": metadata}
        except Exception as e:
            print(f"Error processing file with Textract: {e}")
            return {"extracted_text": "Error processing file with Textract.", "metadata": metadata}
 
    def _process_docx_with_images(self, file_content, metadata):
        """Extracts text and embedded images from DOCX files and processes images for OCR."""
        with open("temp_file.docx", "wb") as temp_file:
            temp_file.write(file_content)
 
        doc = DocxDocument("temp_file.docx")
        text = "\n".join([para.text for para in doc.paragraphs])
 
        # Extract embedded images
        images_dir = "temp_images"
        os.makedirs(images_dir, exist_ok=True)
        extract_images_from_docx("temp_file.docx", images_dir)
       
        image_texts = []
        for image_filename in os.listdir(images_dir):
            if image_filename.endswith((".png", ".jpg", ".jpeg")):
                image_path = os.path.join(images_dir, image_filename)
               
                with open(image_path, "rb") as image_file:
                    image_bytes = image_file.read()
               
                # OCR the image using AWS Textract
                image_text = self._process_pdf_or_image_with_textract(image_bytes, metadata)
                image_texts.append(image_text['extracted_text'])
 
        # Cleanup temporary files and directory
        os.remove("temp_file.docx")
        for image_filename in os.listdir(images_dir):
            os.remove(os.path.join(images_dir, image_filename))
        os.rmdir(images_dir)
 
        return {"extracted_text": text + "\n".join(image_texts), "metadata": metadata}
 
    def _process_pptx(self, file_content, metadata):
        """Extracts text from a PPTX file."""
        temp_pptx_path = "temp_file.pptx"
       
        with open(temp_pptx_path, "wb") as temp_file:
            temp_file.write(file_content)
       
        presentation = Presentation(temp_pptx_path)
        text_runs = []
       
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text.strip())
       
        os.remove(temp_pptx_path)
        return {"extracted_text": "\n".join(text_runs), "metadata": metadata}
 
    def _process_csv(self, file_content, metadata):
        """Extracts text from a CSV file."""
        csv_str = file_content.decode('utf-8')
        df = pd.read_csv(BytesIO(csv_str.encode('utf-8')))
        return {"extracted_text": df.to_string(), "metadata": metadata}
 
    def _process_xlsx(self, file_content, metadata):
        """Extracts text from an XLSX file."""
        xlsx_file = BytesIO(file_content)
        df = pd.read_excel(xlsx_file, sheet_name=None)
        extracted_text = "\n\n".join([f"Sheet: {sheet}\n{df[sheet].to_string()}" for sheet in df])
        return {"extracted_text": extracted_text, "metadata": metadata}
 
    def _process_txt(self, file_content, metadata):
        """Extracts text from a TXT file."""
        return {"extracted_text": file_content.decode('utf-8'), "metadata": metadata}
 
    def process(self, file_content, file_extension, metadata=None):
        """Extracts text based on file type using the appropriate method."""
        metadata = metadata or {}
        if file_extension == ".pdf" or file_extension in [".jpg", ".jpeg", ".png"]:
            return self._process_pdf_or_image_with_textract(file_content, metadata)
        elif file_extension == ".docx":
            return self._process_docx_with_images(file_content, metadata)
        elif file_extension == ".csv":
            return self._process_csv(file_content, metadata)
        elif file_extension == ".xlsx":
            return self._process_xlsx(file_content, metadata)
        elif file_extension == ".txt":
            return self._process_txt(file_content, metadata)
        elif file_extension == ".pptx":
            return self._process_pptx(file_content, metadata)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    def _clean_extracted_text(self, text):
        """Cleans up the extracted text by removing unwanted tags or artifacts."""
        clean_text = re.sub(r'<[^>]+>', '', text)  # Remove any HTML-like tags
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()  # Normalize whitespace
        return clean_text
