import os
import pandas as pd
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader
from docx import Document as DocxDocument
from docx2txt import process as extract_images_from_docx
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
from typing import Dict, Any
from Services.Text_Extraction.iText_Extraction import Text_Service
from datetime import datetime
from pathlib import Path

class AzureService(Text_Service):
    def __init__(self):
        load_dotenv()
        self.endpoint = os.getenv("AZURE_FORM_RECOGNIZER_ENDPOINT")
        self.key = os.getenv("AZURE_FORM_RECOGNIZER_KEY")
        
        # Initialize Document Analysis client
        self.document_analysis_client = DocumentAnalysisClient(
            endpoint=str(self.endpoint),
            credential=AzureKeyCredential(str(self.key))
        )

    def validate_config(self,config: Dict[str, Any]) -> Dict[str, Any]:
        """Validate if necessary configuration is set up properly."""
        required_fields = ["AZURE_FORM_RECOGNIZER_ENDPOINT", "AZURE_FORM_RECOGNIZER_KEY"]
        missing_fields = [field for field in required_fields if not config['config'].get(field)]
        

        if missing_fields:
            return {
                "status": "error",
                "message": f"Missing required environment variables: {', '.join(missing_fields)}",
                "error": True
            }

        # If all checks pass
        return {
            "status": "success",
            "message": "Azure Form Recognizer configuration validated successfully",
            "error": False
        }

    def _extract_text_from_pdf_or_image(self, file_content):
        """Extracts text from PDF or image files using Azure Document Intelligence."""
        temp_file_path = "temp_file.pdf"
        try:
            with open(temp_file_path, "wb") as f:
                f.write(file_content)

            # Use AzureAIDocumentIntelligenceLoader for text extraction
            loader = AzureAIDocumentIntelligenceLoader(
                api_endpoint=self.endpoint,
                api_key=self.key,
                file_path=temp_file_path,
                api_model="prebuilt-read"
            )
            documents = loader.load()

            # Gather extracted text from each page
            extracted_text = "\n".join([doc.page_content.strip() for doc in documents if doc.page_content.strip()])

            # Return response with both text and metadata
            return {
                "status": "success",
                "extracted_text": extracted_text,
                "error": False
            }

        finally:
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)

    def _extract_text_from_docx_with_images(self, file_content):
        """Extracts text and embedded images from a DOCX file."""
        with open("temp_file.docx", "wb") as f:
            f.write(file_content)

        doc = DocxDocument("temp_file.docx")
        text = "\n".join([para.text for para in doc.paragraphs])

        # Extract embedded images
        images_dir = "temp_images"
        os.makedirs(images_dir, exist_ok=True)
        extract_images_from_docx("temp_file.docx", images_dir)

        image_texts = []
        for image_filename in os.listdir(images_dir):
            if image_filename.endswith((".png", ".jpg", ".jpeg")):
                image_path = os.path.join(images_dir, image_filename)

                with open(image_path, "rb") as image_file:
                    image_bytes = image_file.read()

                # OCR the image
                image_text = self._extract_text_from_pdf_or_image(image_bytes)
                image_texts.append(image_text)

        # Cleanup temporary files
        os.remove("temp_file.docx")
        for image_filename in os.listdir(images_dir):
            os.remove(os.path.join(images_dir, image_filename))
        os.rmdir(images_dir)

        return {
            "status": "success",
            "extracted_text": text + "\n".join(image_texts),
            
            "error": False
        }
    
    def _extract_text_from_csv(self, file_content):
        """Extracts text from a CSV file."""
        try:
            csv_str = file_content.decode('utf-8')
            df = pd.read_csv(BytesIO(csv_str.encode('utf-8')))
            return {
            "status": "success",
            "extracted_text": df.to_string(),
            
            "error": False
        }
        except Exception as e:
            return f"Error processing CSV file: {e}"

    def _extract_text_from_xlsx(self, file_content):
        """Extracts text from an XLSX file."""
        try:
            xlsx_file = BytesIO(file_content)
            df = pd.read_excel(xlsx_file, sheet_name=None)
            #return "\n\n".join([f"Sheet: {sheet}\n{df[sheet].to_string()}" for sheet in df])
            return {
            "status": "success",
            "extracted_text":"\n\n".join([f"Sheet: {sheet}\n{df[sheet].to_string()}" for sheet in df]),
            
            "error": False
        }
    
        except Exception as e:
            return f"Error processing XLSX file: {e}"

    def _extract_text_from_txt(self, file_content):
        """Extracts text from a TXT file."""
        try:
           # return file_content.decode('utf-8')
        
          return {
                "status": "success",
                "extracted_text":file_content.decode('utf-8'),
                
                "error": False
            }
        except Exception as e:
            return f"Error processing TXT file: {e}"

    def _extract_text_from_pptx(self, file_content):
        """Extracts text from a PPTX file."""
        temp_pptx_path = "temp_file.pptx"
        try:
            with open(temp_pptx_path, "wb") as f:
                f.write(file_content)

            presentation = Presentation(temp_pptx_path)
            text_runs = []

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text.strip())

            # Return response with both text and metadata
            return {
                "status": "success",
                "extracted_text": "\n".join(text_runs),
                
                "error": False
            }

        finally:
            if os.path.exists(temp_pptx_path):
                os.remove(temp_pptx_path)

    def process(self, file_content, file_extension):
        """Extracts text and metadata based on file type using the appropriate method."""
        try:
            if file_extension in [".pdf", ".jpg", ".jpeg", ".png"]:
                return self._extract_text_from_pdf_or_image(file_content)
            elif file_extension == ".docx":
                return self._extract_text_from_docx_with_images(file_content)
            elif file_extension == ".csv":
                return self._extract_text_from_csv(file_content)
            elif file_extension == ".xlsx":
                return self._extract_text_from_xlsx(file_content)
            elif file_extension == ".txt":
                return self._extract_text_from_txt(file_content)
            elif file_extension == ".pptx":
                return self._extract_text_from_pptx(file_content)
            
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
        except Exception as e:
            return {
                "status": "error",
                "message": f"Error processing the file: {str(e)}",
                "error": True
            }
