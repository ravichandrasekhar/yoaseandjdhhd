import os
import pandas as pd
import pytesseract
from PIL import Image
from pdf2image import convert_from_bytes
from docx import Document as DocxDocument
from docx2txt import process as extract_images_from_docx
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
from Services.Text_Extraction.iText_Extraction import Text_Service
from typing import Dict, Any
 
load_dotenv()
 
class OpensourceService(Text_Service):
    def __init__(self):
        # Set the Poppler and Tesseract paths if provided
        self.poppler_path = os.getenv("POPPLER_PATH")
        self.tesseract_path = os.getenv("TESSERACT_PATH")
        os.environ['PATH'] = self.poppler_path + os.pathsep + os.environ['PATH']
        os.environ['PATH'] = f'{os.path.dirname(self.tesseract_path)};{os.environ["PATH"]}'
 
    def validate_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Validates if necessary configuration is set up correctly."""
        required_fields = ["POPPLER_PATH", "TESSERACT_PATH"]
        missing_fields = [fields for fields in required_fields if not config['config'].get(fields)]
 
        if missing_fields:
            return {
                "status": "error",
                "message": f"Missing required environment variables for Open Source: {', '.join(missing_fields)}",
                "error": True
            }
 
        return {
            "status": "success",
            "message": "Open source configuration validated successfully",
            "error": False
        }
 
    def _extract_text_from_image(self, image_content, metadata):
        """Extracts text from images using Tesseract OCR."""
        image = Image.open(BytesIO(image_content))
        extracted_text = pytesseract.image_to_string(image)
        return {"extracted_text": extracted_text, "metadata": metadata}
 
    def _extract_text_from_pdf(self, file_content, metadata):
        """Extracts text from PDF files using Tesseract OCR."""
        images = convert_from_bytes(file_content, poppler_path=self.poppler_path)
        text = "\n".join([pytesseract.image_to_string(img) for img in images])
        return {"extracted_text": text, "metadata": metadata}
 
    def _extract_text_from_docx_with_images(self, file_content, metadata):
        """Extracts text and images from a DOCX file and processes images for text extraction."""
        temp_file_path = "temp_file.docx"
        images_dir = "temp_images"
 
        try:
            with open(temp_file_path, "wb") as f:
                f.write(file_content)
           
            doc = DocxDocument(temp_file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
 
            # Extract embedded images
            os.makedirs(images_dir, exist_ok=True)
            extract_images_from_docx(temp_file_path, images_dir)
 
            image_texts = []
            for image_filename in os.listdir(images_dir):
                if image_filename.endswith((".png", ".jpg", ".jpeg")):
                    image_path = os.path.join(images_dir, image_filename)
                    with open(image_path, "rb") as image_file:
                        image_text = self._extract_text_from_image(image_file.read(), metadata)
                        image_texts.append(image_text['extracted_text'])
           
            return {"extracted_text": text + "\n".join(image_texts), "metadata": metadata}
        finally:
            # Cleanup temporary files and directories
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            if os.path.exists(images_dir):
                for image_filename in os.listdir(images_dir):
                    os.remove(os.path.join(images_dir, image_filename))
                os.rmdir(images_dir)
 
    def _extract_text_from_csv(self, file_content, metadata):
        """Extracts text from a CSV file."""
        csv_str = file_content.decode('utf-8')
        df = pd.read_csv(BytesIO(csv_str.encode('utf-8')))
        return {"extracted_text": df.to_string(), "metadata": metadata}
 
    def _extract_text_from_xlsx(self, file_content, metadata):
        """Extracts text from an XLSX file."""
        xlsx_file = BytesIO(file_content)
        df = pd.read_excel(xlsx_file, sheet_name=None)
        extracted_text = "\n\n".join([f"Sheet: {sheet}\n{df[sheet].to_string()}" for sheet in df])
        return {"extracted_text": extracted_text, "metadata": metadata}
 
    def _extract_text_from_txt(self, file_content, metadata):
        """Extracts text from a TXT file."""
        return {"extracted_text": file_content.decode('utf-8'), "metadata": metadata}
 
    def _extract_text_from_pptx(self, file_content, metadata):
        """Extracts text from a PPTX file."""
        temp_pptx_path = "temp_file.pptx"
        try:
            with open(temp_pptx_path, "wb") as f:
                f.write(file_content)
           
            presentation = Presentation(temp_pptx_path)
            text_runs = []
 
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text.strip())
           
            return {"extracted_text": "\n".join(text_runs), "metadata": metadata}
        finally:
            if os.path.exists(temp_pptx_path):
                os.remove(temp_pptx_path)
 
    def process(self, file_content, file_extension, metadata=None):
        """Extracts text and metadata based on file type using the appropriate method."""
        metadata = metadata or {}
        if file_extension == ".pdf":
            return self._extract_text_from_pdf(file_content, metadata)
        elif file_extension in [".jpg", ".jpeg", ".png"]:
            return self._extract_text_from_image(file_content, metadata)
        elif file_extension == ".docx":
            return self._extract_text_from_docx_with_images(file_content, metadata)
        elif file_extension == ".csv":
            return self._extract_text_from_csv(file_content, metadata)
        elif file_extension == ".xlsx":
            return self._extract_text_from_xlsx(file_content, metadata)
        elif file_extension == ".txt":
            return self._extract_text_from_txt(file_content, metadata)
        elif file_extension == ".pptx":
            return self._extract_text_from_pptx(file_content, metadata)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
 