from pptx import Presentation
from pydantic import BaseModel
class PresentationExtractor:
    class InputSchema(BaseModel):
        InputSchema:str
    def process_ppt(file_content):
            # Open the PowerPoint presentation
            presentation = Presentation(file_content)
            extracted_text = []

            # Loop through all slides in the presentation
            for slide in presentation.slides:
                for shape in slide.shapes:
                    # Check if the shape has text
                    if hasattr(shape, "text"):
                        extracted_text.append(shape.text)

            return extracted_text